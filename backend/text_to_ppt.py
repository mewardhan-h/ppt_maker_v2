"""
PowerPoint Generator - Converts formatted text to PPT
No terminal prompts - receives all data as parameters
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree
import re
import os


def parse_text(text: str) -> list[dict]:
    """Parse text into slides. Use # for title, - for bullets, | for markdown tables."""
    slides = []
    current_slide = None

    for line in text.strip().splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            if current_slide:
                slides.append(current_slide)
            current_slide = {"title": line[2:].strip(), "bullets": [], "table": None}
        elif line.startswith("- ") and current_slide is not None:
            current_slide["bullets"].append(line[2:].strip())
        elif line.startswith("|") and current_slide is not None:
            # Skip separator rows like |---|---|
            cells = [c.strip() for c in line.strip("|").split("|")]
            if all(re.match(r'^[-:\s]+$', c) for c in cells if c):
                continue
            if current_slide["table"] is None:
                current_slide["table"] = {"rows": []}
            current_slide["table"]["rows"].append(cells)

    if current_slide:
        slides.append(current_slide)

    return slides


# Logo positions and sizes
LOGO_LEFT_WIDTH = Inches(1.54)
LOGO_LEFT_HEIGHT = Inches(1.5)
LOGO_LEFT_LEFT = Inches(0.6)
LOGO_LEFT_TOP = Inches(0.37)

LOGO_RIGHT_WIDTH = Inches(1.95)
LOGO_RIGHT_HEIGHT = Inches(1.5)
LOGO_RIGHT_LEFT = Inches(10.92)
LOGO_RIGHT_TOP = Inches(0.37)


def set_bullet_format(p):
    """Apply hanging indent, black bullet char in text, 1.5 line spacing, left align."""
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.5

    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(0.38))))
    pPr.set('indent', str(-int(Inches(0.38))))

    for tag in ('a:buNone', 'a:buClr', 'a:buChar', 'a:buAutoNum'):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)


def parse_inline_bold(text: str) -> list[tuple]:
    """Split text into (segment, is_bold) pairs based on **markers**."""
    parts = re.split(r'\*\*', text)
    return [(part, i % 2 == 1) for i, part in enumerate(parts) if part]


def set_paragraph_runs(p, bullet_text: str, font_size_pt: int = 24):
    """Clear existing runs and add new ones with inline bold support."""
    p_xml = p._p
    for r in p_xml.findall(qn('a:r')):
        p_xml.remove(r)
    segments = [('\u2022  ', False)] + parse_inline_bold(bullet_text)
    for text, is_bold in segments:
        if not text:
            continue
        r_el = etree.SubElement(p_xml, qn('a:r'))
        rPr = etree.SubElement(r_el, qn('a:rPr'))
        rPr.set('lang', 'en-US')
        rPr.set('sz', str(font_size_pt * 100))
        if is_bold:
            rPr.set('b', '1')
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        t_el = etree.SubElement(r_el, qn('a:t'))
        t_el.text = text


def add_logo(slide, logo_path: str, left, top, width, height):
    """Add logo image to slide"""
    if logo_path and os.path.isfile(logo_path):
        slide.shapes.add_picture(logo_path, left, top, width=width, height=height)


def set_cell_bg(cell, hex_color: str):
    """Set solid background colour on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for existing in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(existing)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_color)


def add_table_to_slide(slide, rows: list, left, top, width, height):
    """Render a markdown table onto the slide."""
    num_rows = len(rows)
    num_cols = max(len(r) for r in rows)
    tbl = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    for row_idx, row in enumerate(rows):
        is_header = row_idx == 0
        for col_idx in range(num_cols):
            cell_text = row[col_idx].strip() if col_idx < len(row) else ""
            cell = tbl.cell(row_idx, col_idx)
            cell.text = cell_text
            tf = cell.text_frame
            tf.margin_left   = Inches(0.05)
            tf.margin_right  = Inches(0.05)
            tf.margin_top    = Inches(0.03)
            tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.bold = is_header
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if is_header else RGBColor(0x00, 0x00, 0x00)
            if is_header:
                set_cell_bg(cell, '1F497D')
            elif row_idx % 2 == 0:
                set_cell_bg(cell, 'DCE6F1')
            else:
                set_cell_bg(cell, 'FFFFFF')


def add_title_slide(prs, topic: str, presenter: str, department: str,
                    logo_left_path: str = None, logo_right_path: str = None):
    """Create the special first slide with topic and presenter details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # Topic text box
    topic_width  = Inches(12.36)
    topic_height = Inches(1.57)
    topic_left   = Inches((13.33 - 12.36) / 2)
    topic_top    = Inches(2.1)
    tb_topic = slide.shapes.add_textbox(topic_left, topic_top, topic_width, topic_height)
    tf_topic = tb_topic.text_frame
    tf_topic.word_wrap = True
    tf_topic.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf_topic.margin_left   = Inches(0.1)
    tf_topic.margin_right  = Inches(0.1)
    tf_topic.margin_top    = Inches(0.05)
    tf_topic.margin_bottom = Inches(0.05)
    p = tf_topic.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = topic
    p.font.name = "Calibri"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Presenter info text box
    info_width  = Inches(7.66)
    info_height = Inches(2.71)
    info_left   = Inches(0.6)
    info_top    = Inches(4.6)
    tb_info = slide.shapes.add_textbox(info_left, info_top, info_width, info_height)
    tf_info = tb_info.text_frame
    tf_info.word_wrap = True
    tf_info.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf_info.margin_left   = Inches(0.1)
    tf_info.margin_right  = Inches(0.1)
    tf_info.margin_top    = Inches(0.05)
    tf_info.margin_bottom = Inches(0.05)

    def add_info_line(tf, is_first, segments):
        p = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        p.space_after  = Pt(4)
        p_xml = p._p
        for r in p_xml.findall(qn('a:r')):
            p_xml.remove(r)
        for text, bold in segments:
            r_el = etree.SubElement(p_xml, qn('a:r'))
            rPr  = etree.SubElement(r_el, qn('a:rPr'))
            rPr.set('lang', 'en-US')
            rPr.set('sz', '2400')
            rPr.set('b', '1' if bold else '0')
            solidFill = etree.SubElement(rPr, qn('a:solidFill'))
            srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', '000000')
            t_el = etree.SubElement(r_el, qn('a:t'))
            t_el.text = text

    name_display = presenter[0].upper() + presenter[1:] if presenter else presenter
    dept_display = department[0].upper() + department[1:] if department else department
    add_info_line(tf_info, True,  [('Presented by: ', False), (name_display, True)])
    add_info_line(tf_info, False, [('Department of ', False), (dept_display, True)])
    add_info_line(tf_info, False, [('Government Villupuram Medical College and Hospital', True)])

    # Add logos
    add_logo(slide, logo_left_path, LOGO_LEFT_LEFT, LOGO_LEFT_TOP, LOGO_LEFT_WIDTH, LOGO_LEFT_HEIGHT)
    add_logo(slide, logo_right_path, LOGO_RIGHT_LEFT, LOGO_RIGHT_TOP, LOGO_RIGHT_WIDTH, LOGO_RIGHT_HEIGHT)


def add_header_topic(slide, title: str):
    """Add slide title centered between the two logos."""
    logo_left_end = LOGO_LEFT_LEFT + LOGO_LEFT_WIDTH
    available_width = LOGO_RIGHT_LEFT - logo_left_end
    txBox = slide.shapes.add_textbox(logo_left_end, LOGO_LEFT_TOP, available_width, LOGO_LEFT_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)


def generate_ppt_with_data(topic: str, presenter: str, department: str, 
                           input_file: str = "basic_text_converted.txt",
                           output_file: str = "presentation.pptx"):
    """
    Generate PowerPoint with provided data (no terminal prompts!)
    
    Args:
        topic: Presentation topic
        presenter: Presenter name
        department: Department name
        input_file: Formatted text file to read
        output_file: Output PPT filename
    
    Returns:
        str: Output file path
    """
    
    print(f"📄 Reading formatted content from: {input_file}")
    
    # Check if input file exists
    if not os.path.isfile(input_file):
        raise FileNotFoundError(f"Input file not found: {input_file}")
    
    # Read formatted content
    with open(input_file, "r", encoding="utf-8") as f:
        content = f.read()
    
    print(f"✅ Read {len(content)} characters")
    
    # Parse slides
    slides = parse_text(content)
    
    if not slides:
        raise ValueError("No slides found in formatted content")
    
    print(f"🎨 Generating {len(slides)} slide(s)...")
    
    # Logo paths
    logo_left_path = "assets/logo_left.jpg"
    logo_right_path = "assets/logo_right.jpg"
    
    # Check logos
    if not os.path.isfile(logo_left_path):
        print(f"⚠️  Warning: Left logo not found: {logo_left_path}")
        logo_left_path = None
    
    if not os.path.isfile(logo_right_path):
        print(f"⚠️  Warning: Right logo not found: {logo_right_path}")
        logo_right_path = None
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    content_layout = prs.slide_layouts[1]
    
    # Title slide (slide 0) - uses provided data
    add_title_slide(prs, topic, presenter, department, logo_left_path, logo_right_path)
    
    # Content slides (slide 1 onwards)
    content_left   = Inches((13.33 - 11.41) / 2)
    content_top    = Inches(2.1)
    content_width  = Inches(11.41)
    content_height = Inches(3.65)
    
    for slide_data in slides[1:]:
        slide = prs.slides.add_slide(content_layout)
        
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)
        
        if slide_data["bullets"]:
            txBox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left   = Inches(0.1)
            tf.margin_right  = Inches(0.1)
            tf.margin_top    = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            for j, bullet in enumerate(slide_data["bullets"]):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                set_bullet_format(p)
                set_paragraph_runs(p, bullet)
        
        if slide_data.get("table"):
            rows = slide_data["table"]["rows"]
            tbl_top = content_top + Inches(len(slide_data["bullets"]) * 0.5) if slide_data["bullets"] else content_top
            tbl_height = content_height - (tbl_top - content_top)
            add_table_to_slide(slide, rows, content_left, tbl_top, content_width, tbl_height)
        
        add_logo(slide, logo_left_path, LOGO_LEFT_LEFT, LOGO_LEFT_TOP, LOGO_LEFT_WIDTH, LOGO_LEFT_HEIGHT)
        add_logo(slide, logo_right_path, LOGO_RIGHT_LEFT, LOGO_RIGHT_TOP, LOGO_RIGHT_WIDTH, LOGO_RIGHT_HEIGHT)
        add_header_topic(slide, slide_data["title"])
    
    # Save presentation
    prs.save(output_file)
    print(f"💾 Saved: {output_file}")
    
    return output_file
